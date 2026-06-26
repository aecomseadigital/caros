"""Caros Office MCP server.

A stdio MCP server exposing full read/create/edit tools for Word (.docx),
Excel (.xlsx), and PowerPoint (.pptx). Bundled with Caros and registered as a
default-enabled `ExtensionConfig::Stdio` extension (see README).

Heavy libraries are imported lazily inside each tool so the process starts fast
and the module imports cleanly even before optional deps are installed.
"""

from __future__ import annotations

from typing import Any

from mcp.server.fastmcp import FastMCP

mcp = FastMCP("caros-office")


# --------------------------------------------------------------------------- #
# Word (.docx)
# --------------------------------------------------------------------------- #
@mcp.tool()
def docx_read(path: str) -> str:
    """Extract all paragraph text from a Word .docx file."""
    from docx import Document

    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)


@mcp.tool()
def docx_create(path: str, paragraphs: list[str], title: str | None = None) -> str:
    """Create a Word .docx file from a list of paragraphs, with an optional title heading."""
    from docx import Document

    doc = Document()
    if title:
        doc.add_heading(title, level=0)
    for text in paragraphs:
        doc.add_paragraph(text)
    doc.save(path)
    return f"Created {path} with {len(paragraphs)} paragraph(s)."


@mcp.tool()
def docx_append(path: str, paragraphs: list[str]) -> str:
    """Append paragraphs to an existing Word .docx file."""
    from docx import Document

    doc = Document(path)
    for text in paragraphs:
        doc.add_paragraph(text)
    doc.save(path)
    return f"Appended {len(paragraphs)} paragraph(s) to {path}."


# --------------------------------------------------------------------------- #
# Excel (.xlsx)
# --------------------------------------------------------------------------- #
@mcp.tool()
def xlsx_read(path: str, sheet: str | None = None) -> list[list[Any]]:
    """Read a worksheet as a list of rows. Defaults to the active sheet."""
    from openpyxl import load_workbook

    wb = load_workbook(path, data_only=True)
    ws = wb[sheet] if sheet else wb.active
    return [list(row) for row in ws.iter_rows(values_only=True)]


@mcp.tool()
def xlsx_write(path: str, rows: list[list[Any]], sheet: str = "Sheet1") -> str:
    """Create (or overwrite) an .xlsx file, writing rows into the named sheet."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = sheet
    for row in rows:
        ws.append(row)
    wb.save(path)
    return f"Wrote {len(rows)} row(s) to sheet '{sheet}' in {path}."


@mcp.tool()
def xlsx_set_cell(path: str, cell: str, value: Any, sheet: str | None = None) -> str:
    """Set a single cell (e.g. 'B2') in an existing .xlsx file."""
    from openpyxl import load_workbook

    wb = load_workbook(path)
    ws = wb[sheet] if sheet else wb.active
    ws[cell] = value
    wb.save(path)
    return f"Set {cell} = {value!r} in {path}."


# --------------------------------------------------------------------------- #
# PowerPoint (.pptx)
# --------------------------------------------------------------------------- #
@mcp.tool()
def pptx_read(path: str) -> list[dict[str, Any]]:
    """Read slides as a list of {index, text} entries (concatenated shape text)."""
    from pptx import Presentation

    prs = Presentation(path)
    slides: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides):
        texts = [
            shape.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text
        ]
        slides.append({"index": i, "text": "\n".join(texts)})
    return slides


@mcp.tool()
def pptx_create(path: str, slides: list[dict[str, Any]]) -> str:
    """Create a .pptx. Each slide is {title, bullets:[...]} rendered as a title+content layout."""
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    for spec in slides:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = str(spec.get("title", ""))
        body = slide.placeholders[1].text_frame
        bullets = spec.get("bullets") or []
        for j, bullet in enumerate(bullets):
            para = body.paragraphs[0] if j == 0 else body.add_paragraph()
            para.text = str(bullet)
    prs.save(path)
    return f"Created {path} with {len(slides)} slide(s)."


if __name__ == "__main__":
    mcp.run()
